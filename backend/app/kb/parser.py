import os
import logging
from typing import Dict, Any, List
from pathlib import Path
import PyPDF2
import docx
from pptx import Presentation
import markdown

from app.core.exceptions import KBParseFailedException

logger = logging.getLogger(__name__)


class DocumentParser:
    """文档解析器"""
    
    @staticmethod
    def parse_file(file_path: str, file_type: str) -> Dict[str, Any]:
        """解析文件"""
        try:
            if file_type.lower() == "pdf":
                return DocumentParser._parse_pdf(file_path)
            elif file_type.lower() == "docx":
                return DocumentParser._parse_docx(file_path)
            elif file_type.lower() == "pptx":
                return DocumentParser._parse_pptx(file_path)
            elif file_type.lower() == "md":
                return DocumentParser._parse_markdown(file_path)
            elif file_type.lower() == "txt":
                return DocumentParser._parse_txt(file_path)
            else:
                raise KBParseFailedException(f"不支持的文件类型: {file_type}")
                
        except Exception as e:
            logger.error(f"文件解析失败 {file_path}: {e}")
            raise KBParseFailedException(f"文件解析失败: {e}")
    
    @staticmethod
    def _parse_pdf(file_path: str) -> Dict[str, Any]:
        """解析PDF文件"""
        content = []
        metadata = {"type": "pdf", "pages": 0}
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            metadata["pages"] = len(pdf_reader.pages)
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                try:
                    text = page.extract_text()
                    if text.strip():
                        content.append({
                            "page": page_num,
                            "text": text.strip(),
                            "section": f"第{page_num}页"
                        })
                except Exception as e:
                    logger.warning(f"PDF页面{page_num}解析失败: {e}")
                    continue
        
        return {
            "content": content,
            "metadata": metadata,
            "raw_text": "\n\n".join([item["text"] for item in content])
        }
    
    @staticmethod
    def _parse_docx(file_path: str) -> Dict[str, Any]:
        """解析Word文档"""
        doc = docx.Document(file_path)
        content = []
        metadata = {"type": "docx", "paragraphs": 0}
        
        current_section = "文档开始"
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
                
            metadata["paragraphs"] += 1
            
            # 检测标题（简单的启发式方法）
            if para.style.name.startswith('Heading') or (
                len(text) < 100 and 
                any(keyword in text for keyword in ['第', '章', '节', '部分', '概述', '介绍'])
            ):
                current_section = text
            
            content.append({
                "paragraph": metadata["paragraphs"],
                "text": text,
                "section": current_section,
                "style": para.style.name
            })
        
        return {
            "content": content,
            "metadata": metadata,
            "raw_text": "\n\n".join([item["text"] for item in content])
        }
    
    @staticmethod
    def _parse_pptx(file_path: str) -> Dict[str, Any]:
        """解析PowerPoint文档"""
        prs = Presentation(file_path)
        content = []
        metadata = {"type": "pptx", "slides": len(prs.slides)}
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_title = f"幻灯片{slide_num}"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_text.append(text)
                    
                    # 第一个文本框通常是标题
                    if len(slide_text) == 1 and len(text) < 100:
                        slide_title = text
            
            if slide_text:
                content.append({
                    "slide": slide_num,
                    "text": "\n".join(slide_text),
                    "section": slide_title
                })
        
        return {
            "content": content,
            "metadata": metadata,
            "raw_text": "\n\n".join([item["text"] for item in content])
        }
    
    @staticmethod
    def _parse_markdown(file_path: str) -> Dict[str, Any]:
        """解析Markdown文件"""
        with open(file_path, 'r', encoding='utf-8') as file:
            md_content = file.read()
        
        # 简单的Markdown解析
        lines = md_content.split('\n')
        content = []
        current_section = "文档开始"
        current_text = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # 检测标题
            if line.startswith('#'):
                # 保存之前的内容
                if current_text:
                    content.append({
                        "text": "\n".join(current_text),
                        "section": current_section
                    })
                    current_text = []
                
                # 更新当前章节
                current_section = line.lstrip('#').strip()
            else:
                current_text.append(line)
        
        # 保存最后的内容
        if current_text:
            content.append({
                "text": "\n".join(current_text),
                "section": current_section
            })
        
        metadata = {"type": "markdown", "sections": len(content)}
        
        return {
            "content": content,
            "metadata": metadata,
            "raw_text": md_content
        }
    
    @staticmethod
    def _parse_txt(file_path: str) -> Dict[str, Any]:
        """解析纯文本文件"""
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
        
        # 简单按段落分割
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        
        content = []
        for i, para in enumerate(paragraphs, 1):
            content.append({
                "paragraph": i,
                "text": para,
                "section": f"段落{i}"
            })
        
        metadata = {"type": "txt", "paragraphs": len(paragraphs)}
        
        return {
            "content": content,
            "metadata": metadata,
            "raw_text": text
        }